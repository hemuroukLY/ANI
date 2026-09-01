"""Document parsing service (US-011 / SPEC §5.1, plan.md §4).

Parses PDF/Word/Excel/PPT/MD/TXT into platform content nodes:

* PDF: lightweight text extraction via PyMuPDF (no model download). Embedded
  images are extracted and uploaded to MinIO as ``[图片: caption](OSS_URL)``
  placeholder nodes. Scanned-page OCR is deferred to a later phase.
* Word/Excel/PPT/MD/TXT: parsed by LlamaIndex DoclingReader.

Post-processing rules applied to all formats:

* Tables → HTML; tables larger than 2048 tokens are split by row groups, each
  preserving the header row.
* Images → uploaded to MinIO and replaced with ``[图片: caption](OSS_URL)``
  placeholder nodes.
"""
from __future__ import annotations

import re
from collections.abc import Iterator
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Literal

from llama_index.readers.docling import DoclingReader

from app.clients.minio_client import ImageUploader

# Markdown table token threshold for row-group splitting (plan.md §4.2).
TABLE_TOKEN_THRESHOLD = 2048
# Rough tokens-per-char heuristic used only for table size estimation.
CHARS_PER_TOKEN = 2

SUPPORTED_EXTS = {".pdf", ".docx", ".xlsx", ".pptx", ".md", ".txt"}
ContentType = Literal["text", "table", "image", "code"]
# Sub-type hint for downstream chunk_service (non-hierarchical context labels).
NodeSubType = Literal["heading", "paragraph", "table", "image", "list", "code"]

_FENCE_RE = re.compile(r"^```", re.MULTILINE)
_CAPTION_RE = re.compile(r"^(?:图|Figure|Fig\.?)\s*\d*[：:.\s]+(?P<cap>.+)$", re.IGNORECASE)
# Markdown ATX heading: ``# Title`` / ``## Title`` … (1-6 levels).
_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*#*$")


@dataclass
class ParsedNode:
    """A logical content node emitted by the parser."""

    content: str
    content_type: ContentType = "text"
    page_number: int | None = None
    metadata: dict[str, Any] = field(default_factory=dict)


def _estimate_tokens(text: str) -> int:
    """Conservative token estimate (chars / 2)."""
    return max(1, len(text) // CHARS_PER_TOKEN)


def _detect_heading(text: str) -> tuple[int, str] | None:
    """Detect a markdown ATX heading in ``text``.

    Returns ``(level, title)`` if the first non-blank line is a heading,
    otherwise ``None``. Level is 1-6 (``#`` = 1, ``######`` = 6).
    """
    for line in text.strip().splitlines():
        if not line.strip():
            continue
        m = _HEADING_RE.match(line)
        if m:
            level = len(m.group(1))
            title = m.group(2).strip()
            return level, title
        # First non-blank line is not a heading.
        return None
    return None


def _is_list_line(line: str) -> bool:
    """Heuristic: line is a markdown list item (``-``, ``*``, ``1.``)."""
    stripped = line.strip()
    if stripped.startswith(("- ", "* ", "+ ")):
        return True
    return bool(re.match(r"^\d+\.\s", stripped))


def _is_list_block(text: str) -> bool:
    """Heuristic: text segment is predominantly a list."""
    lines = [line for line in text.strip().splitlines() if line.strip()]
    if not lines:
        return False
    list_lines = sum(1 for line in lines if _is_list_line(line))
    return list_lines >= len(lines) * 0.5


def _split_by_headings(text: str) -> list[str]:
    """Split text into segments at markdown ATX heading boundaries.

    Each heading line becomes its own segment; subsequent body text until
    the next heading forms the following segment. Returns non-empty text
    segments preserving order.
    """
    segments: list[str] = []
    current: list[str] = []
    for line in text.splitlines():
        if _HEADING_RE.match(line):
            # Flush accumulated body text (if any) before the heading.
            if current:
                segments.append("\n".join(current).strip())
                current = []
            # The heading itself is a standalone segment.
            segments.append(line.strip())
        else:
            current.append(line)
    if current:
        segments.append("\n".join(current).strip())
    return [s for s in segments if s]


def _split_tables_and_text(markdown: str) -> list[tuple[str, str]]:
    """Split markdown into ordered (kind, text) segments; kind ∈ {'text','table'}.

    Docling emits tables as markdown pipe tables (or HTML inside fences). We
    normalise every detected table to an HTML ``<table>`` block and return
    non-table segments verbatim.
    """
    segments: list[tuple[str, str]] = []
    lines = markdown.splitlines()
    i = 0
    n = len(lines)
    buf: list[str] = []

    def flush() -> None:
        if buf:
            segments.append(("text", "\n".join(buf).strip()))
            buf.clear()

    while i < n:
        line = lines[i]
        if _FENCE_RE.match(line):
            fence = line.strip()
            block = [line]
            i += 1
            while i < n and not lines[i].startswith(fence[:3]):
                block.append(lines[i])
                i += 1
            if i < n:
                block.append(lines[i])
                i += 1
            body = "\n".join(block)
            flush()
            if "<table" in body:
                table_html = _extract_html_table(body)
                if table_html:
                    segments.append(("table", table_html))
                else:
                    segments.append(("text", body))
            elif _is_pipe_table(block[1:-1] if len(block) > 2 else []):
                table_html = _pipe_table_to_html("\n".join(block[1:-1]))
                if table_html:
                    segments.append(("table", table_html))
                else:
                    segments.append(("text", body))
            else:
                segments.append(("text", body))
            continue
        if _is_pipe_table([line] + (lines[i + 1 : i + 3] if i + 1 < n else [])):
            block = [line]
            i += 1
            while i < n and _is_pipe_table([lines[i]]):
                block.append(lines[i])
                i += 1
            flush()
            table_html = _pipe_table_to_html("\n".join(block))
            if table_html:
                segments.append(("table", table_html))
            continue
        buf.append(line)
        i += 1
    flush()
    return segments


def _is_pipe_table(lines: list[str]) -> bool:
    """Heuristic: a line is a markdown pipe-table row if it contains pipes."""
    if not lines:
        return False
    row = lines[0]
    if not row.strip().startswith("|"):
        return False
    # Second line should be a separator (|---|---|) when available.
    if len(lines) > 1:
        sep = lines[1].strip()
        if sep.startswith("|") and re.match(r"^\|[\s:|-]+\|\s*$", sep):
            return True
    return "|" in row


def _pipe_table_to_html(table_text: str) -> str:
    """Convert a markdown pipe table to an HTML table (header preserved)."""
    rows = [r for r in table_text.splitlines() if r.strip()]
    if not rows:
        return ""
    # Drop the separator row.
    body = [rows[0]]
    for r in rows[1:]:
        if re.match(r"^\|[\s:|-]+\|\s*$", r.strip()):
            continue
        body.append(r)

    def cells(row: str) -> list[str]:
        return [c.strip() for c in row.strip().strip("|").split("|")]

    out = ["<table>"]
    for idx, r in enumerate(body):
        tag = "th" if idx == 0 else "td"
        out.append(
            "  <tr>" + "".join(f"<{tag}>{c}</{tag}>" for c in cells(r)) + "</tr>"
        )
    out.append("</table>")
    return "\n".join(out)


def _extract_html_table(fenced: str) -> str:
    """Pull the first <table>...</table> block out of a fenced segment."""
    m = re.search(r"<table[\s\S]*?</table>", fenced, re.IGNORECASE)
    return m.group(0) if m else ""


def _split_large_table(html_table: str) -> list[str]:
    """Split a table larger than the token threshold by row groups.

    Each group preserves the header row (first row). A single row larger than
    the threshold becomes its own group (plan.md §4.2).
    """
    if _estimate_tokens(html_table) <= TABLE_TOKEN_THRESHOLD:
        return [html_table]
    header, rows = _decompose_html_table(html_table)
    if not header or not rows:
        # Cannot split without a header row or without data rows to group.
        return [html_table]
    groups: list[str] = []
    current: list[str] = []
    current_tokens = 0
    for row in rows:
        row_tokens = _estimate_tokens(row)
        # If the single row itself exceeds the threshold, emit it alone.
        if row_tokens >= TABLE_TOKEN_THRESHOLD:
            if current:
                groups.append(_compose_table(header, current))
                current = []
                current_tokens = 0
            groups.append(_compose_table(header, [row]))
            continue
        if current and current_tokens + row_tokens > TABLE_TOKEN_THRESHOLD:
            groups.append(_compose_table(header, current))
            current = []
            current_tokens = 0
        current.append(row)
        current_tokens += row_tokens
    if current:
        groups.append(_compose_table(header, current))
    return groups


def _decompose_html_table(html_table: str) -> tuple[str, list[str]]:
    """Return (header_html, [row_html, ...]) from an HTML table."""
    rows = re.findall(r"<tr>[\s\S]*?</tr>", html_table, re.IGNORECASE)
    if not rows:
        return "", []
    return rows[0], rows[1:]


def _compose_table(header: str, rows: list[str]) -> str:
    if not header:
        return "<table>\n" + "\n".join(rows) + "\n</table>"
    return "<table>\n" + header + "\n" + "\n".join(rows) + "\n</table>"


def _html_table_rows(html_table: str) -> tuple[list[str], list[list[str]]]:
    """Extract ``(headers, records)`` from a ``<table>`` HTML block.

    The first row (``<th>``/``<td>`` cells) is treated as the header row and
    returned as a list of plain-text column names. Every following row is a
    data record returned as a list of plain-text cell values (``record[i]``
    aligns with ``headers[i]``).

    All HTML markup is stripped from cell text so callers can safely re-escape
    the values. Returns ``("", [])`` when the table has no header row or no
    data rows (un-decomposable).
    """
    rows = re.findall(r"<tr>[\s\S]*?</tr>", html_table, re.IGNORECASE)
    if not rows:
        return [], []

    # First row cells → column headers.
    headers = _extract_row_cells(rows[0])
    # Every subsequent row is a data record.
    records = [_extract_row_cells(row) for row in rows[1:]]
    records = [r for r in records if r and any(c for c in r)]
    if not headers or not records:
        return [], []
    return headers, records


def _extract_row_cells(row_html: str) -> list[str]:
    """Return the plain-text content of each ``<th>``/``<td>`` cell in a row."""
    cells: list[str] = []
    for cell in re.findall(r"<t[dh][^>]*>([\s\S]*?)</t[dh]>", row_html, re.IGNORECASE):
        text = re.sub(r"<[^>]+>", "", cell)
        text = text.replace("\xa0", " ").strip()
        cells.append(text)
    return cells


def _caption_for(alt: str) -> str | None:
    if not alt:
        return None
    m = _CAPTION_RE.match(alt.strip())
    return m.group("cap") if m else None


def _build_image_placeholder(alt: str, url: str, caption: str | None = None) -> str:
    cap = caption or alt or "图片"
    # Placeholder format ``[图片: caption](OSS_URL)`` (SPEC §5.1 parse) so the
    # frontend can distinguish image nodes from regular markdown images.
    return f"[图片: {cap}]({url})"


class ParseService:
    """Parse documents into platform content nodes using DoclingReader."""

    def __init__(self, uploader: ImageUploader | None = None) -> None:
        # Optional MinIO uploader for extracting and storing embedded images.
        self._uploader = uploader

    def parse(
        self,
        file_path: str,
        object_prefix: str = "",
    ) -> list[ParsedNode]:
        """Parse a single file and return ordered content nodes.

        Args:
            file_path: Local path to the document.
            object_prefix: MinIO key prefix for image uploads
                (``{tenant_id}/{kb_id}/{doc_id}``).

        Returns:
            Ordered list of ParsedNode (text/table/image/code).
        """
        ext = Path(file_path).suffix.lower()
        if ext not in SUPPORTED_EXTS:
            raise ValueError(f"doc.unsupported_type: {ext}")

        if ext == ".pdf":
            # Lightweight PDF path: PyMuPDF text + image extraction, no model
            # download (scanned-page OCR deferred to a later phase).
            return _parse_pdf_lightweight(
                file_path, self._uploader, object_prefix
            )

        if ext in (".md", ".txt"):
            # Plain text / markdown: read directly — DoclingReader mangles
            # pipe tables in TXT/MD files.
            markdown = Path(file_path).read_text(encoding="utf-8")
        else:
            # Word/Excel/PPT: use DoclingReader for text + table extraction.
            # A fresh instance per call avoids state leakage between documents.
            reader = DoclingReader()
            docs = reader.load_data(file_path=file_path)
            markdown = "\n\n".join(d.text for d in docs)

        nodes = _emit_text_table_nodes(markdown)

        # Extract embedded images from Word/Excel/PPT and upload to MinIO.
        if ext in (".docx", ".xlsx", ".pptx") and self._uploader and object_prefix:
            image_nodes = _extract_office_images(
                file_path, ext, self._uploader, object_prefix
            )
            nodes.extend(image_nodes)
        return nodes


def _emit_text_table_nodes(
    markdown: str,
    page_number: int = 1,
) -> list[ParsedNode]:
    """Split markdown into ordered text/table ParsedNodes with context metadata.

    Shared by the MD/TXT/Word/Excel/PPT path and the PDF path so both apply
    the same table-detection, HTML conversion, large-table split, and page
    break logic.

    Each node's ``metadata`` carries non-hierarchical context hints for the
    downstream ``chunk_service``:

    * ``sub_type``: ``heading`` | ``paragraph`` | ``table`` | ``image`` | ``list``
    * ``heading_level``: 1-6 for heading nodes
    * ``section_path``: ``"Chapter > Section > Subsection"`` — breadcrumbs
    * ``table_index`` / ``image_index`` / ``row_count`` / ``is_large_table``
    """
    nodes: list[ParsedNode] = []
    # Breadcrumbs: list of (level, title) for current heading ancestry.
    section_stack: list[tuple[int, str]] = []
    table_index = 0

    def _section_path() -> str:
        return " > ".join(title for _, title in section_stack)

    for kind, segment in _split_tables_and_text(markdown):
        if kind == "table":
            table_index += 1
            is_large = _estimate_tokens(segment) > TABLE_TOKEN_THRESHOLD
            for group in _split_large_table(segment):
                group_rows = group.count("<tr>")
                nodes.append(
                    ParsedNode(
                        content=group,
                        content_type="table",
                        page_number=page_number,
                        metadata={
                            "sub_type": "table",
                            "table_index": table_index,
                            "row_count": group_rows,
                            "is_large_table": is_large,
                            "section_path": _section_path(),
                        },
                    )
                )
        else:
            # Split text segment by heading boundaries so each heading becomes
            # its own node (with sub_type=heading), separate from body text.
            for sub_seg in _split_by_headings(segment):
                for chunk in _split_text_by_page(sub_seg, page_number):
                    page_number = chunk.page_number or page_number
                    # Enrich text node with sub_type and section context.
                    heading = _detect_heading(chunk.content)
                    if heading is not None:
                        level, title = heading
                        # Update breadcrumbs: pop deeper-or-equal levels, push new.
                        while section_stack and section_stack[-1][0] >= level:
                            section_stack.pop()
                        section_stack.append((level, title))
                        chunk.metadata["sub_type"] = "heading"
                        chunk.metadata["heading_level"] = level
                    elif _is_list_block(chunk.content):
                        chunk.metadata["sub_type"] = "list"
                    else:
                        chunk.metadata["sub_type"] = "paragraph"
                    chunk.metadata["section_path"] = _section_path()
                    nodes.append(chunk)

    return nodes


def _extract_office_images(
    file_path: str,
    ext: str,
    uploader: ImageUploader,
    object_prefix: str,
) -> list[ParsedNode]:
    """Extract embedded images from Word/Excel/PPT and upload to MinIO.

    Returns image placeholder nodes (``[图片: caption](OSS_URL)``).
    """
    nodes: list[ParsedNode] = []
    try:
        if ext == ".docx":
            nodes = _extract_docx_images(file_path, uploader, object_prefix)
        elif ext == ".xlsx":
            nodes = _extract_xlsx_images(file_path, uploader, object_prefix)
        elif ext == ".pptx":
            nodes = _extract_pptx_images(file_path, uploader, object_prefix)
    except Exception:  # noqa: BLE001, S110 — best-effort image extraction
        # Image extraction is best-effort; text/table nodes are already emitted.
        pass
    return nodes


def _extract_docx_images(
    file_path: str, uploader: ImageUploader, object_prefix: str
) -> list[ParsedNode]:
    """Extract images from a .docx file via python-docx."""
    import docx

    doc = docx.Document(file_path)
    nodes: list[ParsedNode] = []
    img_idx = 0
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                blob = rel.target_part.blob
                ext_img = rel.target_part.partname.ext
                oss_url = uploader.upload(blob, object_prefix, ext_img.lstrip(".") or "png")
                img_idx += 1
                nodes.append(
                    ParsedNode(
                        content=_build_image_placeholder("", oss_url),
                        content_type="image",
                        page_number=1,
                        metadata={"sub_type": "image", "image_index": img_idx},
                    )
                )
            except Exception:  # noqa: BLE001, S110 — skip unreadable image
                pass
    return nodes


def _extract_xlsx_images(
    file_path: str, uploader: ImageUploader, object_prefix: str
) -> list[ParsedNode]:
    """Extract images from a .xlsx file via openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path)
    nodes: list[ParsedNode] = []
    img_idx = 0
    for ws in wb.worksheets:
        for image in ws._images:
            try:
                data = image._data() if callable(image._data) else image._data
                ext_img = image.format or "png"
                oss_url = uploader.upload(data, object_prefix, ext_img)
                img_idx += 1
                nodes.append(
                    ParsedNode(
                        content=_build_image_placeholder("", oss_url),
                        content_type="image",
                        page_number=1,
                        metadata={"sub_type": "image", "image_index": img_idx},
                    )
                )
            except Exception:  # noqa: BLE001, S110 — skip unreadable image
                pass
    return nodes


def _extract_pptx_images(
    file_path: str, uploader: ImageUploader, object_prefix: str
) -> list[ParsedNode]:
    """Extract images from a .pptx file via python-pptx."""
    from pptx import Presentation

    prs = Presentation(file_path)
    nodes: list[ParsedNode] = []
    img_idx = 0
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    blob = image.blob
                    ext_img = image.ext or "png"
                    oss_url = uploader.upload(blob, object_prefix, ext_img)
                    img_idx += 1
                    nodes.append(
                        ParsedNode(
                            content=_build_image_placeholder("", oss_url),
                            content_type="image",
                            page_number=slide_num,
                            metadata={"sub_type": "image", "image_index": img_idx},
                        )
                    )
                except Exception:  # noqa: BLE001, S110 — skip unreadable image
                    pass
    return nodes


def _pdf_text_with_headings(page: Any) -> str:
    """Extract page text as markdown, converting large-font lines to ATX headings.

    PyMuPDF ``get_text("dict")`` gives per-span font sizes. Lines with a font
    size significantly larger than the page body size (statistical mode) are
    converted to ``#`` / ``##`` / ``###`` headings so downstream
    ``_detect_heading`` and ``_split_by_headings`` can label them.

    Falls back to plain ``get_text("text")`` if dict extraction fails.
    """

    try:
        d = page.get_text("dict")
    except Exception:  # noqa: BLE001 — fall back to plain text extraction
        return page.get_text("text").strip()

    # Collect font sizes to determine the body-text baseline (most common size).
    size_counts: dict[float, int] = {}
    for block in d.get("blocks", []):
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                size = round(span.get("size", 12), 1)
                if size > 0:
                    size_counts[size] = size_counts.get(size, 0) + 1
    if not size_counts:
        return page.get_text("text").strip()

    # Body size = the most frequently occurring font size (mode).
    body_size = max(size_counts, key=size_counts.__getitem__)

    lines: list[str] = []
    for block in d.get("blocks", []):
        for line in block.get("lines", []):
            spans = line.get("spans", [])
            if not spans:
                continue
            text = "".join(s.get("text", "") for s in spans).strip()
            if not text:
                continue
            # Use the max span size in the line as the line's font size.
            line_size = max(round(s.get("size", 12), 1) for s in spans)
            # Map larger-than-body sizes to heading levels.
            if line_size >= body_size * 2.0:
                lines.append(f"# {text}")
            elif line_size >= body_size * 1.5:
                lines.append(f"## {text}")
            elif line_size >= body_size * 1.25:
                lines.append(f"### {text}")
            else:
                lines.append(text)
    return "\n".join(lines)


def _parse_pdf_lightweight(
    file_path: str,
    uploader: ImageUploader | None,
    object_prefix: str,
) -> list[ParsedNode]:
    """Parse a text-based PDF with PyMuPDF — no model download needed.

    Extracts text per page (with font-size-based heading detection) and
    embedded images. Images are uploaded to MinIO (when an uploader is
    provided) and emitted as ``[图片: caption](OSS_URL)`` placeholder nodes.
    Scanned-page OCR is deferred to a later phase.
    """
    import fitz  # PyMuPDF

    nodes: list[ParsedNode] = []
    img_idx = 0
    with fitz.open(file_path) as doc:
        for page_num in range(len(doc)):
            page = doc[page_num]
            page_number = page_num + 1

            # 1. Extract text with heading detection → text/table nodes.
            text = _pdf_text_with_headings(page)
            if text:
                nodes.extend(_emit_text_table_nodes(text, page_number))

            # 2. Extract embedded images and upload to MinIO.
            if uploader and object_prefix:
                for img_info in page.get_images(full=True):
                    xref = img_info[0]
                    try:
                        base_image = doc.extract_image(xref)
                        img_bytes = base_image["image"]
                        img_ext = base_image.get("ext", "png")
                        oss_url = uploader.upload(
                            img_bytes, object_prefix, img_ext
                        )
                        img_idx += 1
                        placeholder = _build_image_placeholder("", oss_url)
                        nodes.append(
                            ParsedNode(
                                content=placeholder,
                                content_type="image",
                                page_number=page_number,
                                metadata={
                                    "sub_type": "image",
                                    "image_index": img_idx,
                                },
                            )
                        )
                    except Exception:  # noqa: BLE001, S110 — skip unreadable image
                        # Skip images that cannot be extracted.
                        pass
    return nodes


def _split_text_by_page(text: str, start_page: int) -> Iterator[ParsedNode]:
    """Yield text nodes, advancing page_number on Docling page markers."""
    parts = re.split(r"<!--\s*page[:\s]*(\d+)\s*-->", text)
    page = start_page
    # re.split with a group yields [pre, num, post, num, post, ...]
    idx = 0
    while idx < len(parts):
        seg = parts[idx]
        if seg.strip():
            yield ParsedNode(content=seg.strip(), content_type="text", page_number=page)
        if idx + 1 < len(parts):
            try:
                page = int(parts[idx + 1])
            except ValueError:
                page += 1
            idx += 2
        else:
            idx += 1
