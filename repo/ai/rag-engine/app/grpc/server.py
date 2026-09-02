"""rag-engine gRPC server — Parse/Embed/Generate/GenerateStream RPCs
(RAG-REFACTOR-STEP-2 / Plan §2.2-§2.5).

Implements the gRPC service defined in ``rag.proto``:

* ``Parse`` — stateless document parse + chunk (download_url → chunks + images).
* ``Embed`` — stateless embedding (texts → vectors_flat).
* ``Generate`` — stateless LLM completion (question + context + history → answer).
* ``GenerateStream`` — streaming variant of Generate.

The stateless RPCs do NOT depend on LlamaIndex (Plan §2.3-§2.4).

Error mapping (SPEC §4.4)::

    INVALID_ARGUMENT  — question empty / top_k out of range
    NOT_FOUND         — kb_id collection missing (handled by retriever)
    UNAVAILABLE       — vLLM unavailable
    DEADLINE_EXCEEDED — LLM timeout

The server runs on a background thread with its own ``grpc.aio`` event loop
so it cooperates with FastAPI's event loop when both are started in the same
process (``main.py`` starts the gRPC server in a separate thread to keep
the FastAPI loop responsive).
"""
from __future__ import annotations

import asyncio
import json
import logging
import os
import threading
import uuid
from collections.abc import AsyncGenerator
from pathlib import Path
from typing import Any

import grpc

from app.core.config import settings
from app.grpc import rag_pb2 as rag_pb
from app.grpc import rag_pb2_grpc as rag_grpc

logger = logging.getLogger(__name__)

# gRPC server default bind address. Override via settings.grpc_bind_addr.
DEFAULT_GRPC_BIND = "[::]:50052"

# Default child chunk size for Parse RPC when request.chunk_size is 0.
DEFAULT_PARSE_CHUNK_SIZE = 1024

# P3: httpx download timeout (seconds) — prevents indefinite hang on slow
# download URLs. gRPC client deadline will still apply as an outer bound.
DOWNLOAD_TIMEOUT_SECONDS = 120.0


class RagEngineServicer(rag_grpc.RagEngineServicer):
    """Implements the stateless Parse/Embed/Generate/GenerateStream RPCs.

    The servicer lazily initializes stateless service instances (ParseService,
    EmbedRPCService, GenerateRPCService) on first use and reuses them across
    requests to avoid leaking httpx connection pools.
    """

    def __init__(self) -> None:
        # P6: Reuse stateless service instances across requests to avoid
        # leaking httpx connection pools on every RPC call.
        self._lock = threading.Lock()
        self._parse_svc: Any = None
        self._embed_svc: Any = None
        self._generate_svc: Any = None

    # ── Parse RPC (Plan §2.2) ──────────────────────────────────────────────

    async def Parse(
        self,
        request: rag_pb.ParseRequest,
        context: grpc.aio.ServicerContext,
    ) -> rag_pb.ParseResponse:
        """Stateless document parse + chunk (Plan §2.2).

        Downloads from ``download_url`` to a temp file, reuses
        ``ParseService.parse`` + ``ChunkService.chunk``, extracts image bytes
        (PDF/Office), and returns ``ParsedChunk`` list (parents + children +
        images). Summary is NOT generated here (kb-service orchestrates it
        via Generate RPC).
        """
        if not request.download_url:
            await context.abort(
                grpc.StatusCode.INVALID_ARGUMENT, "download_url is required"
            )
            return  # type: ignore[return-value]
        if not request.file_name:
            await context.abort(
                grpc.StatusCode.INVALID_ARGUMENT, "file_name is required"
            )
            return  # type: ignore[return-value]

        import tempfile

        import httpx

        from app.services.chunk_service import ChunkService
        from app.services.parse_service import ParseService

        chunk_size = request.chunk_size if request.chunk_size else DEFAULT_PARSE_CHUNK_SIZE
        # Determine file extension for temp file suffix.
        ext = request.file_type or os.path.splitext(request.file_name)[1]
        if ext and not ext.startswith("."):
            ext = "." + ext

        # 1. Download to temp file (httpx streaming, P3: with timeout).
        with tempfile.NamedTemporaryFile(suffix=ext or ".tmp", delete=False) as f:
            temp_path = f.name
            download_error: Exception | None = None
            try:
                async with (
                    httpx.AsyncClient(timeout=DOWNLOAD_TIMEOUT_SECONDS) as client,
                    client.stream("GET", request.download_url) as resp,
                ):
                    resp.raise_for_status()
                    async for chunk in resp.aiter_bytes():
                        f.write(chunk)
            except Exception as exc:  # noqa: BLE001
                download_error = exc

        if download_error is not None:
            _safe_unlink(temp_path)
            await context.abort(
                grpc.StatusCode.INTERNAL, f"download failed: {download_error}"
            )
            return  # type: ignore[return-value]

        try:
            # 2. Parse + chunk (offload blocking calls to thread).
            # P6: Reuse ParseService instance (stateless for RPC path).
            if self._parse_svc is None:
                with self._lock:
                    if self._parse_svc is None:  # DCL: second check
                        self._parse_svc = ParseService(uploader=None)
            nodes = await asyncio.to_thread(self._parse_svc.parse, temp_path, "")

            # For md files with image:// references, strip the raw references
            # from text nodes — actual image bytes are extracted separately
            # in step 3 and image links are appended by kb-service.
            if (request.file_type or ext or "").lstrip(".").lower() in ("md", "txt"):
                import re as _re2
                _strip_re = _re2.compile(r"!\[[^\]]*\]\(image://[a-f0-9\-]+\)")
                for node in nodes:
                    if node.content_type == "text":
                        node.content = _strip_re.sub("", node.content).strip()

            chunk_svc = ChunkService(parent_chunk_size=chunk_size)
            parents, children = await asyncio.to_thread(chunk_svc.chunk, nodes)

            # 3. Extract image bytes (independent of parse_service uploader).
            file_type = (request.file_type or ext or "").lstrip(".").lower()
            # For md files, pass Core API base URL + auth headers so image://
            # references can be downloaded.
            _core_base = settings.ani_gateway_url.rstrip("/")
            if not _core_base.endswith("/api/v1"):
                _core_base = _core_base.rstrip("/") + "/api/v1"
            _core_hdrs: dict[str, str] = {}
            _dev_tenant = os.environ.get("ANI_DEV_TENANT_ID", "")
            if _dev_tenant:
                _core_hdrs["X-Dev-Tenant-ID"] = _dev_tenant
            _core_token = os.environ.get("ANI_CORE_API_TOKEN", "")
            if _core_token:
                _core_hdrs["Authorization"] = f"Bearer {_core_token}"
            image_chunks = await asyncio.to_thread(
                _extract_image_bytes, temp_path, file_type,
                _core_base, _core_hdrs,
            )

            # 4. Assemble ParsedChunk list.
            chunks = _build_parse_chunks(parents, children, image_chunks)
            return rag_pb.ParseResponse(chunks=chunks)
        except ValueError as exc:
            await context.abort(grpc.StatusCode.INVALID_ARGUMENT, str(exc))
            return  # type: ignore[return-value]
        except Exception as exc:
            logger.exception("rag-engine Parse failed")
            await context.abort(grpc.StatusCode.INTERNAL, str(exc))
            return  # type: ignore[return-value]
        finally:
            _safe_unlink(temp_path)

    # ── Embed RPC (Plan §2.3) ──────────────────────────────────────────────

    async def Embed(
        self,
        request: rag_pb.EmbedRequest,
        context: grpc.aio.ServicerContext,
    ) -> rag_pb.EmbedResponse:
        """Stateless embedding (Plan §2.3).

        Calls ``OpenAICompatibleEmbedding.get_text_embedding_batch`` directly
        (no LlamaIndex). Returns flattened vectors + dimension + count.
        """
        if not request.texts:
            return rag_pb.EmbedResponse(vectors_flat=[], dimension=0, count=0)

        from app.services.embed_rpc_service import EmbedRPCService

        if self._embed_svc is None:
            with self._lock:
                if self._embed_svc is None:  # DCL: second check
                    self._embed_svc = EmbedRPCService()
        svc = self._embed_svc
        try:
            vectors, dim = await asyncio.to_thread(svc.embed, list(request.texts))
        except Exception as exc:
            logger.exception("rag-engine Embed failed")
            await context.abort(grpc.StatusCode.INTERNAL, str(exc))
            return  # type: ignore[return-value]

        # Flatten: vectors_flat[i * dim + j] = j-th dim of i-th text.
        vectors_flat: list[float] = []
        for vec in vectors:
            vectors_flat.extend(vec)
        return rag_pb.EmbedResponse(
            vectors_flat=vectors_flat,
            dimension=dim,
            count=len(vectors),
        )

    # ── Generate RPC (Plan §2.4) ───────────────────────────────────────────

    async def Generate(
        self,
        request: rag_pb.GenerateRequest,
        context: grpc.aio.ServicerContext,
    ) -> rag_pb.GenerateResponse:
        """Stateless LLM generation (Plan §2.4).

        Pure Python ``openai`` SDK + CompactAndRefine context truncation
        reproduction. History includes current-turn user message; question is
        appended as final USER (reproduces ``{query_str}`` template).
        """
        from app.services.generate_rpc_service import GenerateRPCService

        if self._generate_svc is None:
            with self._lock:
                self._generate_svc = GenerateRPCService()
        svc = self._generate_svc
        # Convert proto SourceChunk + ChatMessage to dicts.
        context_dicts = [
            {
                "chunk_id": c.chunk_id,
                "doc_id": c.doc_id,
                "file_name": c.file_name,
                "page": c.page,
                "content": c.content,
                "score": c.score,
            }
            for c in request.context
        ]
        history_dicts = [
            {"role": m.role, "content": m.content}
            for m in request.history
        ]
        max_tokens = request.max_tokens if request.max_tokens else 2048

        try:
            result = await asyncio.to_thread(
                svc.generate,
                request.question,
                request.session_id,
                context_dicts,
                history_dicts,
                request.inference_service_name,
                max_tokens,
            )
        except TimeoutError as exc:
            logger.warning("rag-engine Generate LLM timeout: %s", exc)
            await context.abort(grpc.StatusCode.DEADLINE_EXCEEDED, "LLM timed out")
            return  # type: ignore[return-value]
        except RuntimeError as exc:
            logger.warning("rag-engine Generate backend unavailable: %s", exc)
            await context.abort(grpc.StatusCode.UNAVAILABLE, str(exc))
            return  # type: ignore[return-value]
        except Exception as exc:
            logger.exception("rag-engine Generate failed")
            await context.abort(grpc.StatusCode.INTERNAL, str(exc))
            return  # type: ignore[return-value]

        return rag_pb.GenerateResponse(
            answer=result["answer"],
            input_tokens=result["input_tokens"],
            output_tokens=result["output_tokens"],
            session_id=result["session_id"],
        )

    # ── GenerateStream RPC (Plan §2.4) ─────────────────────────────────────

    async def GenerateStream(
        self,
        request: rag_pb.GenerateRequest,
        context: grpc.aio.ServicerContext,
    ) -> AsyncGenerator[rag_pb.GenerateToken, None]:
        """Streaming LLM generation (Plan §2.4).

        Yields token events then a final done event with usage. Uses
        ``stream_options={"include_usage": True}`` so the last chunk carries
        usage. Event sequence: token* → done.

        P4: The synchronous ``generate_stream`` generator is iterated in a
        worker thread to avoid blocking the aio event loop.
        """
        from app.services.generate_rpc_service import GenerateRPCService

        if self._generate_svc is None:
            with self._lock:
                if self._generate_svc is None:  # DCL: second check
                    self._generate_svc = GenerateRPCService()
        svc = self._generate_svc
        context_dicts = [
            {
                "chunk_id": c.chunk_id,
                "doc_id": c.doc_id,
                "file_name": c.file_name,
                "page": c.page,
                "content": c.content,
                "score": c.score,
            }
            for c in request.context
        ]
        history_dicts = [
            {"role": m.role, "content": m.content}
            for m in request.history
        ]
        max_tokens = request.max_tokens if request.max_tokens else 2048

        # P4: The synchronous ``generate_stream`` generator is iterated in a
        # worker thread to avoid blocking the aio event loop. Tokens are
        # relayed back via an asyncio.Queue using call_soon_threadsafe so the
        # event loop is never blocked on queue.get.
        loop = asyncio.get_running_loop()
        token_queue: asyncio.Queue[dict | None] = asyncio.Queue()
        generation_error: Exception | None = None

        def _run_stream() -> None:
            nonlocal generation_error
            try:
                for token in svc.generate_stream(
                    request.question,
                    request.session_id,
                    context_dicts,
                    history_dicts,
                    request.inference_service_name,
                    max_tokens,
                ):
                    loop.call_soon_threadsafe(token_queue.put_nowait, token)
            except Exception as exc:  # noqa: BLE001
                generation_error = exc
            finally:
                loop.call_soon_threadsafe(token_queue.put_nowait, None)  # sentinel

        thread = threading.Thread(target=_run_stream, daemon=True)
        thread.start()

        try:
            while True:
                item = await token_queue.get()
                if item is None:
                    break
                yield rag_pb.GenerateToken(
                    content=item.get("content", ""),
                    done=item.get("done", False),
                    input_tokens=item.get("input_tokens", 0),
                    output_tokens=item.get("output_tokens", 0),
                )
            # Worker thread ended — check if it was due to an error.
            if generation_error is not None:
                if isinstance(generation_error, TimeoutError):
                    await context.abort(
                        grpc.StatusCode.DEADLINE_EXCEEDED, "LLM timed out"
                    )
                    return  # type: ignore[return-value]
                elif isinstance(generation_error, RuntimeError):
                    await context.abort(
                        grpc.StatusCode.UNAVAILABLE, str(generation_error)
                    )
                    return  # type: ignore[return-value]
                else:
                    await context.abort(
                        grpc.StatusCode.INTERNAL, str(generation_error)
                    )
                    return  # type: ignore[return-value]
        except TimeoutError as exc:
            logger.warning("rag-engine GenerateStream LLM timeout: %s", exc)
            await context.abort(grpc.StatusCode.DEADLINE_EXCEEDED, "LLM timed out")
            return  # type: ignore[return-value]
        except RuntimeError as exc:
            logger.warning("rag-engine GenerateStream backend unavailable: %s", exc)
            await context.abort(grpc.StatusCode.UNAVAILABLE, str(exc))
            return  # type: ignore[return-value]
        except Exception as exc:
            logger.exception("rag-engine GenerateStream failed")
            await context.abort(grpc.StatusCode.INTERNAL, str(exc))
            return  # type: ignore[return-value]


# ── Parse RPC helpers (Plan §2.2) ────────────────────────────────────────────


def _safe_unlink(path: str) -> None:
    """Best-effort unlink that ignores Windows file-lock errors."""
    try:
        os.unlink(path)
    except OSError:
        pass


def _extract_image_bytes(
    file_path: str,
    file_type: str,
    core_api_base_url: str = "",
    core_api_headers: dict | None = None,
) -> list[dict]:
    """Extract embedded image bytes from a document (Plan §2.2).

    Independent implementation — does NOT depend on parse_service's uploader.
    Returns ``[{"image_bytes", "image_format", "placeholder"}, ...]``.

    PDF: ``doc.extract_image(xref)`` (matches old _parse_pdf_lightweight).
    Office: python-docx / openpyxl / python-pptx dedicated APIs.
    MD: scans ``![alt](image://{object_id})`` references and downloads
        images from Core API via ``core_api_base_url``.
    """
    images: list[dict] = []
    ft = (file_type or "").lower()
    if ft == "pdf":
        import fitz  # PyMuPDF

        with fitz.open(file_path) as doc:
            for page in doc:
                for img_info in page.get_images(full=True):
                    xref = img_info[0]
                    try:
                        base_image = doc.extract_image(xref)
                        images.append({
                            "image_bytes": base_image["image"],
                            "image_format": base_image.get("ext", "png"),
                            "placeholder": "[图片](placeholder)",
                        })
                    except Exception:  # noqa: BLE001, S110 — skip unreadable
                        pass
    elif ft == "docx":
        import docx

        doc = docx.Document(file_path)
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    images.append({
                        "image_bytes": rel.target_part.blob,
                        "image_format": rel.target_part.partname.ext.lstrip(".") or "png",
                        "placeholder": "[图片](placeholder)",
                    })
                except Exception:  # noqa: BLE001, S110 — skip unreadable
                    pass
    elif ft == "xlsx":
        from openpyxl import load_workbook

        wb = load_workbook(file_path)
        for ws in wb.worksheets:
            for img in ws._images:
                try:
                    # P5: img._data() returns a BytesIO object, not bytes.
                    # Read the bytes from it. If it's already bytes, use as-is.
                    raw = img._data() if callable(img._data) else img._data
                    if hasattr(raw, "read"):
                        raw = raw.read()
                    images.append({
                        "image_bytes": raw,
                        "image_format": img.format or "png",
                        "placeholder": "[图片](placeholder)",
                    })
                except Exception:  # noqa: BLE001, S110 — skip unreadable
                    pass
        wb.close()
    elif ft == "pptx":
        from pptx import Presentation

        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        images.append({
                            "image_bytes": image.blob,
                            "image_format": image.ext.lstrip(".") or "png",
                            "placeholder": "[图片](placeholder)",
                        })
                    except Exception:  # noqa: BLE001, S110 — skip unreadable
                        pass
    elif ft == "md":
        # Scan markdown for ``image://{object_id}`` references and download
        # images from Core API.
        import re as _re

        import httpx

        _md_img_re = _re.compile(r"!\[([^\]]*)\]\(image://([a-f0-9\-]+)\)")
        try:
            md_text = Path(file_path).read_text(encoding="utf-8")
        except Exception:  # noqa: BLE001
            md_text = ""
        if md_text and core_api_base_url:
            hdrs = core_api_headers or {}
            for m in _md_img_re.finditer(md_text):
                object_id = m.group(2).strip()
                try:
                    with httpx.Client(timeout=30.0) as client:
                        resp = client.get(
                            f"{core_api_base_url}/objects/{object_id}/download",
                            headers=hdrs,
                        )
                        if resp.status_code != 200:
                            logger.warning("md image download failed (%d) for %s", resp.status_code, object_id)
                            continue
                        download_url = resp.json().get("download_url", "")
                        if not download_url:
                            continue
                        img_resp = client.get(download_url, timeout=60.0)
                        if img_resp.status_code != 200:
                            continue
                        images.append({
                            "image_bytes": img_resp.content,
                            "image_format": "png",
                            "placeholder": "[图片](placeholder)",
                        })
                except Exception as exc:  # noqa: BLE001 — best-effort
                    logger.warning("md image extraction failed for %s: %s", object_id, exc)
    return images


def _build_parse_chunks(
    parents: list[Any],
    children: list[Any],
    image_chunks: list[dict],
) -> list[rag_pb.ParsedChunk]:
    """Assemble ParsedChunk list: parents + children + images (Plan §2.2).

    Summary (doc_summary) is NOT generated here — kb-service orchestrates it
    via Generate RPC. Child chunks carry ``parent_content`` +
    ``parent_chunk_id`` (for _return_parent_and_dedup dedup key + parent
    backfill).

    Args:
        parents: List of ``ParentChunk`` dataclass instances.
        children: List of ``ChildChunk`` dataclass instances.
        image_chunks: List of dicts with keys ``image_bytes``,
            ``image_format``, ``placeholder``.
    """
    chunks = []
    for p in parents:
        chunks.append(rag_pb.ParsedChunk(
            chunk_id=p.chunk_id,
            content=p.content,
            content_type=p.content_type,
            page_number=p.page_number or 0,
            parent_content="",
            parent_chunk_id="",
            chunk_type="parent",
            metadata_json=json.dumps(p.metadata, default=str),
        ))
    for c in children:
        chunks.append(rag_pb.ParsedChunk(
            chunk_id=c.chunk_id,
            content=c.content,
            content_type=c.content_type,
            page_number=c.page_number or 0,
            parent_content=c.parent_content or "",
            parent_chunk_id=c.parent_chunk_id or "",
            chunk_type="child",
            metadata_json=json.dumps(c.metadata, default=str),
        ))
    for img in image_chunks:
        chunks.append(rag_pb.ParsedChunk(
            chunk_id=str(uuid.uuid4()),
            content=img["placeholder"],
            content_type="image",
            page_number=0,
            parent_content="",
            parent_chunk_id="",
            chunk_type="image",
            metadata_json="{}",
            image_bytes=img["image_bytes"],
            image_format=img["image_format"],
        ))
    return chunks


class GrpcServer:
    """Manages the gRPC server lifecycle.

    The server runs on a background thread with its own ``grpc.aio`` event
    loop so it can coexist with FastAPI. ``start()`` is non-blocking;
    ``stop()`` gracefully drains via ``run_coroutine_threadsafe`` (#1).
    """

    def __init__(
        self,
        *,
        bind_addr: str | None = None,
        servicer: RagEngineServicer | None = None,
    ) -> None:
        self._bind_addr = bind_addr or str(getattr(settings, "grpc_bind_addr", DEFAULT_GRPC_BIND))
        self._servicer = servicer or RagEngineServicer()
        self._server: grpc.aio.Server | None = None
        self._thread: threading.Thread | None = None
        self._started = threading.Event()
        # #1: event loop owned by the background thread; stop() uses
        # run_coroutine_threadsafe to schedule stop() on the correct loop.
        self._loop: asyncio.AbstractEventLoop | None = None

    @property
    def bind_addr(self) -> str:
        return self._bind_addr

    @property
    def servicer(self) -> RagEngineServicer:
        return self._servicer

    def start(self) -> None:
        """Start the gRPC server on a background thread (non-blocking)."""
        if self._server is not None:
            return
        self._thread = threading.Thread(target=self._run, name="rag-grpc-server", daemon=True)
        self._thread.start()
        self._started.wait(timeout=10.0)

    def _run(self) -> None:
        # #1: capture this thread's event loop so stop() can schedule on it.
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        self._loop = loop
        try:
            loop.run_until_complete(self._serve_async())
        finally:
            loop.close()
            self._loop = None

    async def _serve_async(self) -> None:
        self._server = grpc.aio.server()
        rag_grpc.add_RagEngineServicer_to_server(self._servicer, self._server)
        self._server.add_insecure_port(self._bind_addr)
        await self._server.start()
        self._started.set()
        logger.info("rag-engine gRPC server listening on %s", self._bind_addr)
        await self._server.wait_for_termination()

    def stop(self, grace: float = 5.0) -> None:
        """Stop the gRPC server gracefully.

        #1: Uses ``asyncio.run_coroutine_threadsafe`` to schedule
        ``server.stop()`` on the background thread's event loop (where the
        aio server was created), rather than the caller's loop.
        """
        if self._server is None:
            return
        loop = self._loop
        if loop is not None and not loop.is_closed():
            try:
                future = asyncio.run_coroutine_threadsafe(
                    self._server.stop(grace=grace), loop
                )
                # Wait for the stop coroutine to complete (bounded by grace).
                future.result(timeout=grace + 2)
            except Exception as exc:  # noqa: BLE001
                logger.warning("rag-engine gRPC server stop failed: %s", exc)
        if self._thread is not None:
            self._thread.join(timeout=grace + 2)
            self._thread = None
        self._server = None
        self._loop = None


def serve(
    *,
    bind_addr: str | None = None,
    block: bool = True,
) -> GrpcServer:
    """Start a gRPC server (convenience entrypoint).

    Args:
        bind_addr: ``host:port`` to bind. Defaults to settings.grpc_bind_addr
            or ``[::]:50052``.
        block: When ``True`` blocks the calling thread until the server
            terminates (process entrypoint). When ``False`` starts in the
            background and returns the :class:`GrpcServer`.
    """
    server = GrpcServer(bind_addr=bind_addr, servicer=RagEngineServicer())
    if block:
        # #14: unified path — reuse GrpcServer._serve_async instead of a
        # duplicate _serve_blocking function.
        asyncio.run(server._serve_async())
    else:
        server.start()
    return server


if __name__ == "__main__":  # pragma: no cover
    logging.basicConfig(level=logging.INFO)
    serve(block=True)
